import streamlit as st
import anthropic
from google.oauth2.service_account import Credentials
from googleapiclient.discovery import build
import pandas as pd
from dotenv import load_dotenv
import os
import io
import PyPDF2
import docx
import openpyxl
from pptx import Presentation

load_dotenv()

SCOPES = [
    "https://www.googleapis.com/auth/spreadsheets",
    "https://www.googleapis.com/auth/drive"
]

def get_credentials():
    try:
        creds_dict = dict(st.secrets["gcp_service_account"])
        creds_dict["private_key"] = creds_dict["private_key"].replace("\\n", "\n")
        return Credentials.from_service_account_info(creds_dict, scopes=SCOPES)
    except:
        return Credentials.from_service_account_file("credentials.json", scopes=SCOPES)

def parse_input_url(url):
    url = url.strip()
    if "folders/" in url:
        return "folder", url.split("folders/")[1].split("?")[0]
    elif "spreadsheets/d/" in url:
        return "sheet", url.split("spreadsheets/d/")[1].split("/")[0]
    elif "presentation/d/" in url:
        return "slides", url.split("presentation/d/")[1].split("/")[0]
    elif "document/d/" in url:
        return "doc", url.split("document/d/")[1].split("/")[0]
    else:
        return "folder", url.strip()

def scan_folder(folder_id):
    creds = get_credentials()
    service = build("drive", "v3", credentials=creds)
    all_files = []

    def scan_recursive(fid):
        results = service.files().list(
            q=f"'{fid}' in parents and trashed=false",
            fields="files(id, name, mimeType)"
        ).execute()
        for f in results.get("files", []):
            if f["mimeType"] == "application/vnd.google-apps.folder":
                scan_recursive(f["id"])
            else:
                all_files.append(f)

    scan_recursive(folder_id)
    return all_files

def export_as_text(file_id):
    creds = get_credentials()
    service = build("drive", "v3", credentials=creds)
    request = service.files().export_media(fileId=file_id, mimeType="text/plain")
    return request.execute().decode("utf-8")

def export_as_csv(file_id):
    creds = get_credentials()
    service = build("drive", "v3", credentials=creds)
    request = service.files().export_media(fileId=file_id, mimeType="text/csv")
    return request.execute().decode("utf-8")

def download_file(file_id):
    creds = get_credentials()
    service = build("drive", "v3", credentials=creds)
    request = service.files().get_media(fileId=file_id)
    return request.execute()

def read_pdf(file_id):
    content = download_file(file_id)
    reader = PyPDF2.PdfReader(io.BytesIO(content))
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
    return text.strip()

def read_docx(file_id):
    content = download_file(file_id)
    doc = docx.Document(io.BytesIO(content))
    return "\n".join([para.text for para in doc.paragraphs if para.text])

def read_xlsx(file_id):
    content = download_file(file_id)
    wb = openpyxl.load_workbook(io.BytesIO(content))
    text = ""
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        text += f"\n--- Sheet: {sheet} ---\n"
        for row in ws.iter_rows(values_only=True):
            row_text = "\t".join([str(c) if c is not None else "" for c in row])
            text += row_text + "\n"
    return text

def read_pptx(file_id):
    content = download_file(file_id)
    prs = Presentation(io.BytesIO(content))
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()

def read_file(f):
    mime = f["mimeType"]
    fid = f["id"]

    if mime == "application/vnd.google-apps.spreadsheet":
        return "[Google Sheet]\n" + export_as_csv(fid)

    elif mime == "application/vnd.google-apps.document":
        return "[Google Doc]\n" + export_as_text(fid)

    elif mime == "application/vnd.google-apps.presentation":
        return "[Google Slides]\n" + export_as_text(fid)

    elif mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        text = read_pptx(fid)
        return "[PPTX]\n" + text if text else None

    elif mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return "[DOCX]\n" + read_docx(fid)

    elif mime == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        return "[XLSX]\n" + read_xlsx(fid)

    elif mime == "application/pdf":
        text = read_pdf(fid)
        if not text:
            return None
        return "[PDF]\n" + text

    return None

def read_all_files(folder_id):
    files = scan_folder(folder_id)
    all_data = {}
    errors = []
    skipped = []

    for f in files:
        try:
            content = read_file(f)
            if content is None:
                skipped.append(f"{f['name']} (ข้ามเพราะอ่านไม่ได้หรือเป็น PDF รูปภาพ)")
            else:
                all_data[f["name"]] = content
        except Exception as e:
            errors.append(f"{f['name']}: {str(e)}")

    return all_data, files, errors, skipped

def ask_claude(question, all_data):
    try:
        api_key = st.secrets["ANTHROPIC_API_KEY"]
    except:
        api_key = os.getenv("ANTHROPIC_API_KEY")

    client = anthropic.Anthropic(api_key=api_key)
    context = "\n\n".join([f"=== {name} ===\n{content}"
                           for name, content in all_data.items()])

    message = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=1024,
        messages=[{
            "role": "user",
            "content": f"""คุณเป็น AI assistant ที่ช่วยวิเคราะห์ข้อมูล ตอบในภาษาเดียวกับคำถาม

ข้อมูลจาก Google Drive:
{context}

คำถาม: {question}"""
        }]
    )
    return message.content[0].text

st.title("🤖 Chatbot จาก Google Drive")
st.markdown("---")

folder_input = st.text_input(
    "📁 วาง Google Drive URL (โฟลเดอร์, Sheet, Doc, Slides, PDF, PPTX, DOCX, XLSX):",
    placeholder="https://drive.google.com/drive/folders/xxxx"
)

if folder_input:
    input_type, file_id = parse_input_url(folder_input)

    with st.spinner("กำลังโหลดข้อมูล..."):
        try:
            if input_type == "folder":
                all_data, files, errors, skipped = read_all_files(file_id)
                st.success(f"พบไฟล์ทั้งหมด {len(files)} ไฟล์ — อ่านได้ {len(all_data)} ไฟล์")

                with st.expander("ดูไฟล์ทั้งหมด"):
                    for f in files:
                        icon = "📊" if "spreadsheet" in f["mimeType"] else "📄"
                        st.write(f"{icon} {f['name']}")

                if skipped:
                    with st.expander(f"⏭ ข้ามไป {len(skipped)} ไฟล์"):
                        for s in skipped:
                            st.write(f"- {s}")

                if errors:
                    with st.expander(f"❌ Error {len(errors)} ไฟล์"):
                        for e in errors:
                            st.error(e)

            elif input_type == "sheet":
                content = export_as_csv(file_id)
                all_data = {"Google Sheet": f"[Google Sheet]\n{content}"}
                errors = []
                st.success("โหลด Google Sheet สำเร็จ")

            elif input_type == "slides":
                content = export_as_text(file_id)
                all_data = {"Google Slides": f"[Google Slides]\n{content}"}
                errors = []
                st.success("โหลด Google Slides สำเร็จ")

            elif input_type == "doc":
                content = export_as_text(file_id)
                all_data = {"Google Doc": f"[Google Doc]\n{content}"}
                errors = []
                st.success("โหลด Google Doc สำเร็จ")

            if "messages" not in st.session_state:
                st.session_state.messages = []
            if "file_id" not in st.session_state or st.session_state.file_id != file_id:
                st.session_state.messages = []
                st.session_state.file_id = file_id

            for msg in st.session_state.messages:
                with st.chat_message(msg["role"]):
                    st.write(msg["content"])

            if question := st.chat_input("ถามอะไรก็ได้เกี่ยวกับข้อมูลนี้..."):
                st.session_state.messages.append({"role": "user", "content": question})
                with st.chat_message("user"):
                    st.write(question)

                with st.chat_message("assistant"):
                    with st.spinner("กำลังคิด..."):
                        answer = ask_claude(question, all_data)
                        st.write(answer)
                        st.session_state.messages.append({"role": "assistant", "content": answer})

        except Exception as e:
            st.error(f"เกิดข้อผิดพลาด: {e}")
            st.exception(e)
else:
    st.info("กรุณาใส่ Google Drive URL เพื่อเริ่มต้น")